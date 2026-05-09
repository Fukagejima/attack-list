"""
PowerPoint (.pptx) → マスキング済み Word (.docx) 変換
プロンプト仕様に従い、全スライドを以下の構成で Word 出力:
  見出し1: スライド[番号]: [タイトル]
  見出し2: 本文 / 表 / 図形・テキストボックス / ノート
"""
from __future__ import annotations
import io

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document
from docx.shared import Pt
from docx.oxml.ns import qn

from .patterns import apply_masks


def _get_slide_title(slide) -> str:
    """スライドのタイトルプレースホルダーからテキストを取得"""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            return ph.text_frame.text.strip()
    return ''


def _collect_shapes(slide):
    """スライド内のシェイプを種別ごとに分類"""
    body_texts = []
    table_shapes = []
    other_texts = []

    for shape in slide.shapes:
        # テーブル
        if shape.has_table:
            table_shapes.append(shape.table)
            continue

        if not shape.has_text_frame:
            continue

        tf = shape.text_frame
        paras = [p.text for p in tf.paragraphs if p.text.strip()]
        if not paras:
            continue

        # プレースホルダー（タイトル以外）= 本文
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
            if shape.placeholder_format.idx != 0:
                body_texts.extend(paras)
        else:
            # テキストボックス・図形
            other_texts.extend(paras)

    return body_texts, table_shapes, other_texts


def _add_table_to_doc(doc: Document, table, opts: dict) -> int:
    """pptx テーブルを Word テーブルとして追加"""
    count = 0
    n_rows = len(table.rows)
    n_cols = len(table.columns)

    if n_rows == 0 or n_cols == 0:
        return 0

    wtable = doc.add_table(rows=n_rows, cols=n_cols)
    wtable.style = 'Table Grid'

    for r_idx, row in enumerate(table.rows):
        for c_idx, cell in enumerate(row.cells):
            text = cell.text_frame.text if cell.text_frame else ''
            masked, cnt = apply_masks(text, opts)
            wtable.rows[r_idx].cells[c_idx].text = masked
            count += cnt

    return count


def mask_ppt(file_bytes: bytes, opts: dict) -> tuple[bytes, int]:
    """
    PPT ファイルをマスクして Word バイト列で返す。
    return: (word_bytes, total_mask_count)
    """
    prs = Presentation(io.BytesIO(file_bytes))
    doc = Document()

    # デフォルトフォント設定
    style = doc.styles['Normal']
    style.font.name = 'Noto Sans JP'
    style.font.size = Pt(10.5)

    total = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        raw_title = _get_slide_title(slide) or f'（タイトルなし）'
        masked_title, cnt = apply_masks(raw_title, opts)
        total += cnt

        # ── スライド見出し ──
        doc.add_heading(f'スライド{slide_num}: {masked_title}', level=1)

        body_texts, table_shapes, other_texts = _collect_shapes(slide)

        # ── 本文 ──
        if body_texts:
            doc.add_heading('本文', level=2)
            for text in body_texts:
                masked, cnt = apply_masks(text, opts)
                total += cnt
                doc.add_paragraph(masked, style='List Bullet')

        # ── 表 ──
        for table in table_shapes:
            doc.add_heading('表', level=2)
            total += _add_table_to_doc(doc, table, opts)

        # ── 図形・テキストボックス ──
        if other_texts:
            doc.add_heading('図形・テキストボックス', level=2)
            for text in other_texts:
                masked, cnt = apply_masks(text, opts)
                total += cnt
                doc.add_paragraph(masked, style='List Bullet')

        # ── ノート ──
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_text = notes_tf.text.strip() if notes_tf else ''
            if notes_text:
                doc.add_heading('ノート', level=2)
                masked, cnt = apply_masks(notes_text, opts)
                total += cnt
                doc.add_paragraph(masked)

        # スライド間セパレータ
        doc.add_paragraph('─' * 40)

    out = io.BytesIO()
    doc.save(out)
    return out.getvalue(), total
